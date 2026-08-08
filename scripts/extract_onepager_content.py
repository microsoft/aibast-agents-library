#!/usr/bin/env python3
"""
Extract the approved AI BAST one-pager slides into structured catalog content.

The source decks use a consistent one-slide layout. This script reads the
PowerPoint shapes directly so customer challenges, agent actions, outcomes,
personas, and Microsoft products remain verbatim and traceable.

Usage:
    python scripts/extract_onepager_content.py \
        --source ~/Desktop/aibast_bible/onepagers_pptx \
        --out state/onepager_content.json
"""

import argparse
import hashlib
import json
import re
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parent.parent
DEFAULT_OUT = ROOT / "state" / "onepager_content.json"
EMU_PER_INCH = 914400


def iter_shapes(shapes):
    """Yield leaf shapes, including children nested inside PowerPoint groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def text_lines(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    return [line.strip() for line in shape.text.splitlines() if line.strip()]


def shape_rows(slide):
    rows = []
    for shape in iter_shapes(slide.shapes):
        lines = text_lines(shape)
        if not lines:
            continue
        rows.append({
            "top": shape.top / EMU_PER_INCH,
            "left": shape.left / EMU_PER_INCH,
            "lines": lines,
        })
    return rows


def closest(rows, predicate, target_top):
    matches = [row for row in rows if predicate(row)]
    if not matches:
        return None
    return min(matches, key=lambda row: abs(row["top"] - target_top))


def split_block(lines):
    return {
        "intro": lines[0] if lines else "",
        "items": lines[1:] if len(lines) > 1 else [],
    }


def metadata_table(slide):
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        if len(table.rows) < 2:
            continue
        headers = [cell.text.strip().lower() for cell in table.rows[0].cells]
        values = [
            [line.strip() for line in cell.text.splitlines() if line.strip()]
            for cell in table.rows[1].cells
        ]
        return {
            headers[index]: values[index]
            for index in range(min(len(headers), len(values)))
        }
    return {}


def extract_onepager(path):
    presentation = Presentation(path)
    slide = presentation.slides[0]
    rows = shape_rows(slide)
    metadata = metadata_table(slide)

    title_row = closest(
        rows,
        lambda row: 0.5 <= row["top"] <= 2.2
        and row["left"] < 0.6
        and len(" ".join(row["lines"])) < 120,
        1.5,
    )
    summary_row = closest(
        rows,
        lambda row: 1.8 <= row["top"] <= 3.4
        and row["left"] < 0.6
        and len(" ".join(row["lines"])) >= 40,
        2.5,
    )
    scenario_row = next(
        (
            row for row in rows
            if any("agent in action example scenario" in line.lower()
                   for line in row["lines"])
        ),
        None,
    )

    columns = []
    for left_min, left_max in ((4.0, 7.0), (7.0, 10.4), (10.4, 14.0)):
        columns.append(closest(
            rows,
            lambda row, lo=left_min, hi=left_max:
                row["top"] >= 4.5
                and lo <= row["left"] < hi
                and len(row["lines"]) >= 2,
            5.0,
        ))

    missing = [
        label for label, value in (
            ("title", title_row),
            ("executive summary", summary_row),
            ("scenario", scenario_row),
            ("situation column", columns[0]),
            ("agent actions column", columns[1]),
            ("outcomes column", columns[2]),
        )
        if value is None
    ]
    if missing:
        raise ValueError(f"{path.name}: missing {', '.join(missing)}")

    opportunity_rows = [
        row for row in rows
        if 3.0 <= row["top"] <= 7.0
        and row["left"] < 4.0
        and len(" ".join(row["lines"])) > 15
        and "situation" not in " ".join(row["lines"]).lower()
    ]
    opportunity_rows.sort(key=lambda row: row["top"])

    modules = []
    opportunities = []
    for row in opportunity_rows:
        lines = row["lines"]
        if len(lines) >= 2 and lines[0].lower().endswith("agent"):
            modules.append({
                "name": lines[0],
                "description": " ".join(lines[1:]),
            })
        else:
            opportunities.append(" ".join(lines))

    scenario_text = " ".join(scenario_row["lines"])
    scenario_name = re.sub(
        r"^.*?agent in action example scenario\s*:\s*",
        "",
        scenario_text,
        flags=re.IGNORECASE,
    ).strip()

    raw = path.read_bytes()
    return {
        "source_file": path.name,
        "source_sha256": hashlib.sha256(raw).hexdigest(),
        "slide_count": len(presentation.slides),
        "title": " ".join(title_row["lines"]),
        "executive_summary": " ".join(summary_row["lines"]),
        "industries": metadata.get("industries", []),
        "personas": metadata.get("personas", []),
        "agent_requirements": metadata.get("agent requirements", []),
        "featured_tools": metadata.get("featured tools", []),
        "scenario_name": scenario_name,
        "customer_challenge": split_block(columns[0]["lines"]),
        "agent_actions": split_block(columns[1]["lines"]),
        "business_outcomes": split_block(columns[2]["lines"]),
        "opportunity_statements": opportunities,
        "modules": modules,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--source",
        type=Path,
        default=Path.home() / "Desktop" / "aibast_bible" / "onepagers_pptx",
        help="Directory containing approved AI BAST one-pager PPTX files",
    )
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()

    source = args.source.expanduser().resolve()
    files = sorted(source.glob("*.pptx"))
    if not files:
        parser.error(f"no .pptx files found in {source}")

    records = {}
    failures = []
    for path in files:
        try:
            records[path.name] = extract_onepager(path)
        except Exception as error:
            failures.append(f"{path.name}: {error}")

    if failures:
        print("\n".join(failures))
        return 1

    doc = {
        "schema": "aibast-onepager-content/1.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source": "Approved AI BAST Agents Library one-pager PowerPoint export",
        "stats": {
            "onepagers": len(records),
            "with_modules": sum(bool(row["modules"]) for row in records.values()),
        },
        "onepagers": records,
    }
    args.out.write_text(json.dumps(doc, indent=2) + "\n", encoding="utf-8")
    print(f"[OK] Extracted {len(records)} one-pagers to {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
